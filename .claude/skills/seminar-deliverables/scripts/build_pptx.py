#!/usr/bin/env python3
"""세미나 발표덱 빌더. deck_spec.json -> .pptx

spec 스키마:
{
  "title": "세미나 제목",
  "subtitle": "thesis 한 문장",
  "meta": "발표자 · 날짜",
  "slides": [
    {"layout": "section", "title": "섹션 제목"},
    {"layout": "bullets", "title": "한 메시지 문장", "bullets": ["20자 내외", ...],
     "notes": "발표자 노트 (근거 타임코드 포함)"},
    {"layout": "quote", "quote": "인용문", "attrib": "발표자 (00:12:05)"},
    {"layout": "stat", "value": "42%", "label": "지표명", "basis": "2024년 대비"}
  ]
}
layout 미지정 시 bullets. 알 수 없는 layout은 bullets로 처리하고 stderr에 경고.
"""
import argparse, json, sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

INK, INK2, ACCENT, BG = RGBColor(0x16,0x18,0x1D), RGBColor(0x4A,0x4F,0x5A), RGBColor(0x2F,0x6D,0xF6), RGBColor(0xFB,0xFB,0xFC)
FONT = "Pretendard"
W, H = Inches(13.333), Inches(7.5)


def _blank(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = BG
    return s


def _text(slide, x, y, w, h, size, color, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.font.size, r.font.bold, r.font.name = Pt(size), bold, FONT
    r.font.color.rgb = color
    return tf, r


def _accent_bar(slide, y=Inches(1.5), w=Inches(1.1)):
    bar = slide.shapes.add_shape(1, Inches(0.9), y, w, Pt(5))
    bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT; bar.line.fill.background()


def _notes(slide, text):
    if text:
        slide.notes_slide.notes_text_frame.text = text


def cover(prs, spec):
    s = _blank(prs)
    _accent_bar(s, Inches(2.4))
    _text(s, Inches(0.9), Inches(2.8), Inches(11.5), Inches(1.6), 44, INK, bold=True)[1].text = spec.get("title", "")
    if spec.get("subtitle"):
        _text(s, Inches(0.9), Inches(4.5), Inches(10.5), Inches(1.2), 20, INK2)[1].text = spec["subtitle"]
    if spec.get("meta"):
        _text(s, Inches(0.9), Inches(6.2), Inches(10.5), Inches(0.6), 14, INK2)[1].text = spec["meta"]


def section(prs, sl):
    s = _blank(prs)
    _accent_bar(s, Inches(3.2))
    _text(s, Inches(0.9), Inches(3.6), Inches(11.5), Inches(1.4), 34, INK, bold=True)[1].text = sl.get("title", "")
    _notes(s, sl.get("notes"))


def bullets(prs, sl):
    s = _blank(prs)
    _accent_bar(s)
    _text(s, Inches(0.9), Inches(1.9), Inches(11.5), Inches(1.3), 30, INK, bold=True)[1].text = sl.get("title", "")
    items = sl.get("bullets") or []
    if items:
        box = s.shapes.add_textbox(Inches(0.9), Inches(3.4), Inches(11.5), Inches(3.4))
        tf = box.text_frame; tf.word_wrap = True
        for i, b in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.space_after = Pt(14)
            r = p.add_run(); r.text = "· " + str(b)
            r.font.size, r.font.name, r.font.color.rgb = Pt(19), FONT, INK2
    _notes(s, sl.get("notes"))


def quote(prs, sl):
    s = _blank(prs)
    bar = s.shapes.add_shape(1, Inches(0.9), Inches(2.4), Pt(5), Inches(2.6))
    bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT; bar.line.fill.background()
    _text(s, Inches(1.5), Inches(2.5), Inches(10.5), Inches(2.4), 26, INK)[1].text = "“" + sl.get("quote", "") + "”"
    if sl.get("attrib"):
        _text(s, Inches(1.5), Inches(5.2), Inches(10.5), Inches(0.6), 15, INK2)[1].text = "— " + sl["attrib"]
    _notes(s, sl.get("notes"))


def stat(prs, sl):
    s = _blank(prs)
    _text(s, Inches(0.9), Inches(2.3), Inches(11.5), Inches(2.0), 88, ACCENT, bold=True, align=PP_ALIGN.CENTER)[1].text = str(sl.get("value", ""))
    _text(s, Inches(0.9), Inches(4.5), Inches(11.5), Inches(0.8), 22, INK, align=PP_ALIGN.CENTER)[1].text = sl.get("label", "")
    if sl.get("basis"):
        _text(s, Inches(0.9), Inches(5.4), Inches(11.5), Inches(0.6), 13, INK2, align=PP_ALIGN.CENTER)[1].text = sl["basis"]
    _notes(s, sl.get("notes"))


LAYOUTS = {"section": section, "bullets": bullets, "quote": quote, "stat": stat}


def build(spec, out):
    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H
    cover(prs, spec)
    for i, sl in enumerate(spec.get("slides", [])):
        name = sl.get("layout", "bullets")
        fn = LAYOUTS.get(name)
        if fn is None:
            print(f"warn: slide {i} unknown layout {name!r}, using bullets", file=sys.stderr)
            fn = bullets
        fn(prs, sl)
    prs.save(out)
    return len(prs.slides._sldIdLst)


def demo():
    import tempfile, os
    spec = {"title": "T", "subtitle": "S", "meta": "M", "slides": [
        {"layout": "section", "title": "sec"},
        {"layout": "bullets", "title": "b", "bullets": ["1", "2"], "notes": "n"},
        {"layout": "quote", "quote": "q", "attrib": "a"},
        {"layout": "stat", "value": "42%", "label": "l", "basis": "b"},
        {"title": "no-layout defaults to bullets"},
        {"layout": "nope", "title": "unknown falls back"},
    ]}
    p = os.path.join(tempfile.mkdtemp(), "d.pptx")
    build(spec, p)
    prs = Presentation(p)
    n = len(prs.slides._sldIdLst)
    assert n == 7, f"expected 7 slides (cover + 6), got {n}"
    assert prs.slides[2].notes_slide.notes_text_frame.text == "n"
    assert os.path.getsize(p) > 20000
    print("demo ok:", p)


if __name__ == "__main__":
    ap = argparse.ArgumentParser()
    ap.add_argument("--spec"); ap.add_argument("--out")
    ap.add_argument("--demo", action="store_true")
    a = ap.parse_args()
    if a.demo:
        demo()
    else:
        if not a.spec or not a.out:
            ap.error("--spec and --out required (or --demo)")
        with open(a.spec, encoding="utf-8") as f:
            build(json.load(f), a.out)
        print(a.out)
